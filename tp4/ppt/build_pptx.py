"""Construye TP4_presentacion.pptx — tema oscuro, gráficos grandes."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "TP4_presentacion.pptx")

# ─── Paleta tema oscuro ──────────────────────────────────────────────────────
BG        = RGBColor(0x0d, 0x11, 0x17)
PANEL     = RGBColor(0x16, 0x1b, 0x22)
BORDER    = RGBColor(0x30, 0x36, 0x3d)
ACCENT_T  = RGBColor(0x5e, 0xea, 0xd4)
ACCENT_P  = RGBColor(0xc0, 0x84, 0xfc)
ACCENT_O  = RGBColor(0xfb, 0xbf, 0x24)
TEXT      = RGBColor(0xf9, 0xfa, 0xfb)
MUTED     = RGBColor(0x9c, 0xa3, 0xaf)
RED       = RGBColor(0xf8, 0x71, 0x71)
GREEN     = RGBColor(0x4a, 0xde, 0x80)

# ─── Layout constants ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# Margins (tighter than before)
MARGIN_X       = Inches(0.45)
CONTENT_W      = SW - 2 * MARGIN_X         # 12.43"
HEADER_TOP     = Inches(0.30)
TITLE_H        = Inches(0.55)
DIVIDER_TOP    = Inches(1.05)              # justo bajo el título
CONTENT_TOP    = Inches(1.20)
FOOTER_TOP     = Inches(7.10)
CONTENT_H      = FOOTER_TOP - CONTENT_TOP  # 5.90"

# Side panel preset (para slides image + texto)
SIDE_PANEL_W   = Inches(3.20)
IMG_PLUS_GAP   = Inches(0.20)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.45),
                   color=ACCENT_T):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_panel(slide, left, top, width, height, *,
              fill=PANEL, border=BORDER, border_w=0.75, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if radius:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(border_w)
    sh.shadow.inherit = False
    return sh


def add_header(slide, title, *, section=None, accent=ACCENT_T):
    add_accent_bar(slide, MARGIN_X, HEADER_TOP + Inches(0.05),
                   width=Inches(0.10), height=Inches(0.55), color=accent)
    label_x = MARGIN_X + Inches(0.22)
    if section:
        add_text(slide, label_x, HEADER_TOP, Inches(10), Inches(0.26),
                 section, size=10, bold=True, color=accent, font="Calibri")
        add_text(slide, label_x, HEADER_TOP + Inches(0.26), Inches(12), TITLE_H,
                 title, size=24, bold=True, color=TEXT, font="Calibri")
    else:
        add_text(slide, label_x, HEADER_TOP + Inches(0.05), Inches(12), TITLE_H,
                 title, size=26, bold=True, color=TEXT, font="Calibri")
    # divider
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  MARGIN_X, DIVIDER_TOP,
                                  CONTENT_W, Emu(9525))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER
    line.line.fill.background()


def add_footer(slide, idx, total):
    add_text(slide, MARGIN_X, FOOTER_TOP + Inches(0.05), Inches(9), Inches(0.3),
             "TP4 — Aprendizaje No Supervisado · SIA 2026 · Grupo G02",
             size=9, color=MUTED)
    add_text(slide, SW - Inches(1.4) - MARGIN_X, FOOTER_TOP + Inches(0.05),
             Inches(1.4), Inches(0.3),
             f"{idx} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def fit_image(slide, path, left, top, max_w, max_h):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    cx = left + (max_w - w) // 2
    cy = top + (max_h - h) // 2
    return slide.shapes.add_picture(path, cx, cy, width=w, height=h)


def img_panel(slide, path, left, top, panel_w, panel_h, *, caption=None,
              caption_size=10, padding=Inches(0.10)):
    add_panel(slide, left, top, panel_w, panel_h, radius=True)
    cap_h = Inches(0.30) if caption else 0
    fit_image(slide, path,
              left + padding, top + padding,
              panel_w - 2 * padding,
              panel_h - 2 * padding - cap_h)
    if caption:
        add_text(slide, left + padding,
                 top + panel_h - padding - cap_h + Inches(0.02),
                 panel_w - 2 * padding, cap_h,
                 caption, size=caption_size, color=MUTED,
                 align=PP_ALIGN.CENTER)


def full_image(slide, path, *, caption=None):
    """Imagen ocupando todo el área de contenido."""
    img_panel(slide, path, MARGIN_X, CONTENT_TOP, CONTENT_W, CONTENT_H,
              caption=caption, padding=Inches(0.12))


def image_with_side(slide, path, side_render, *, caption=None):
    """Imagen grande a izquierda + panel lateral renderizado por side_render(left, top, w, h)."""
    img_w = CONTENT_W - SIDE_PANEL_W - IMG_PLUS_GAP
    img_panel(slide, path, MARGIN_X, CONTENT_TOP, img_w, CONTENT_H,
              caption=caption, padding=Inches(0.12))
    side_left = MARGIN_X + img_w + IMG_PLUS_GAP
    side_render(side_left, CONTENT_TOP, SIDE_PANEL_W, CONTENT_H)


def two_images(slide, paths_captions):
    """Dos imágenes lado a lado ocupando el área de contenido."""
    p1, c1 = paths_captions[0]
    p2, c2 = paths_captions[1]
    gap = Inches(0.20)
    w = (CONTENT_W - gap) / 2
    img_panel(slide, p1, MARGIN_X, CONTENT_TOP, w, CONTENT_H,
              caption=c1, padding=Inches(0.12))
    img_panel(slide, p2, MARGIN_X + w + gap, CONTENT_TOP, w, CONTENT_H,
              caption=c2, padding=Inches(0.12))


def three_images(slide, paths_captions):
    gap = Inches(0.18)
    w = (CONTENT_W - 2 * gap) / 3
    for i, (p, c) in enumerate(paths_captions):
        left = MARGIN_X + (w + gap) * i
        img_panel(slide, p, left, CONTENT_TOP, w, CONTENT_H,
                  caption=c, padding=Inches(0.10))


def grid_2x2(slide, paths_captions):
    gap = Inches(0.15)
    w = (CONTENT_W - gap) / 2
    h = (CONTENT_H - gap) / 2
    for i, (p, c) in enumerate(paths_captions):
        r, col = divmod(i, 2)
        left = MARGIN_X + (w + gap) * col
        top = CONTENT_TOP + (h + gap) * r
        img_panel(slide, p, left, top, w, h, caption=c, padding=Inches(0.10))


def side_panel_text(side_left, side_top, side_w, side_h, *,
                    heading, items, accent=ACCENT_T, item_size=12):
    """Helper para crear un panel lateral con título + bullets."""
    def render(L, T, W, H):
        add_panel(None, L, T, W, H, radius=True)  # placeholder; will be replaced below
    # Esta función SOLO devuelve el closure que pinta en el slide actual.
    raise NotImplementedError("Use side_panel() inline.")


def side_panel(slide, L, T, W, H, *, heading, items, accent=ACCENT_T,
               item_size=12):
    add_panel(slide, L, T, W, H, radius=True)
    pad = Inches(0.22)
    add_text(slide, L + pad, T + pad, W - 2 * pad, Inches(0.45),
             heading, size=14, bold=True, color=accent)
    bullet_list(slide, L + pad, T + pad + Inches(0.55),
                W - 2 * pad, H - pad - Inches(0.55), items,
                size=item_size, bullet_color=accent)


def kpi_box(slide, left, top, width, height, label, value, *, value_color=ACCENT_T,
            value_size=28, label_size=10):
    add_panel(slide, left, top, width, height, radius=True)
    add_text(slide, left, top + Inches(0.18), width, Inches(0.32),
             label.upper(), size=label_size, color=MUTED,
             align=PP_ALIGN.CENTER, bold=True)
    add_text(slide, left, top + Inches(0.55), width, Inches(0.9),
             value, size=value_size, bold=True, color=value_color,
             align=PP_ALIGN.CENTER)


def bullet_list(slide, left, top, width, height, items, *, size=14, color=TEXT,
                bullet_color=ACCENT_T):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run_b = p.add_run()
        run_b.text = "▸ "
        run_b.font.size = Pt(size)
        run_b.font.color.rgb = bullet_color
        run_b.font.bold = True
        run_b.font.name = "Calibri"
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(size)
        run_t.font.color.rgb = color
        run_t.font.name = "Calibri"
    return tb


def add_table(slide, left, top, width, height, data, *,
              header_bg=ACCENT_T, header_fg=BG, body_fg=TEXT, font_size=12,
              col_widths=None):
    rows = len(data); cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = table_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = cw
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg if ri == 0 else PANEL
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.bold = (ri == 0)
            run.font.color.rgb = header_fg if ri == 0 else body_fg
            run.font.name = "Calibri"
    return tbl


# ─── Slide queue ────────────────────────────────────────────────────────────
queue = []


def slide_portada(s):
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SH)
    box.fill.solid(); box.fill.fore_color.rgb = ACCENT_T; box.line.fill.background()
    add_text(s, Inches(0.9), Inches(2.05), Inches(11), Inches(0.55),
             "SIA · TRABAJO PRÁCTICO 4", size=14, bold=True, color=ACCENT_T)
    add_text(s, Inches(0.9), Inches(2.6), Inches(12), Inches(1.4),
             "Aprendizaje No Supervisado", size=52, bold=True, color=TEXT)
    add_text(s, Inches(0.9), Inches(3.9), Inches(12), Inches(0.6),
             "Kohonen · Oja · Hopfield", size=24, color=ACCENT_P)
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.9), Inches(4.85),
                             Inches(3.5), Emu(15000))
    div.fill.solid(); div.fill.fore_color.rgb = BORDER
    div.line.fill.background()
    add_text(s, Inches(0.9), Inches(5.0), Inches(12), Inches(0.4),
             "Grupo G02 · 1er Cuatrimestre 2026", size=16, color=MUTED)
    add_text(s, Inches(0.9), Inches(5.45), Inches(12), Inches(0.4),
             "Dataset: europe.csv · 28 países · 7 features", size=14, color=MUTED)
    add_text(s, Inches(0.9), Inches(6.6), Inches(12), Inches(0.4),
             "commit 28217db", size=11, color=MUTED)
queue.append(slide_portada)


def slide_roadmap(s):
    add_header(s, "Contenido")
    card_w = Inches(4.05); card_h = Inches(5.4)
    gap = Inches(0.20); start = MARGIN_X; top = CONTENT_TOP + Inches(0.10)
    titles = ["Ejercicio 1.1\nKohonen (SOM)",
              "Ejercicio 1.2\nRegla de Oja",
              "Ejercicio 2\nHopfield"]
    descs = [
        ["Country map", "U-Matrix", "Hit map", "Comparación vs PCA"],
        ["Loadings PC1", "Ranking de países", "Convergencia",
         "Comparación vs sklearn"],
        ["Recall 4 letras (GRTV)", "Estado espúreo",
         "Ortogonalidad + crosstalk", "Capacidad + escalado adaptativo"],
    ]
    colors = [ACCENT_T, ACCENT_P, ACCENT_O]
    for i in range(3):
        L = start + (card_w + gap) * i
        add_panel(s, L, top, card_w, card_h, radius=True)
        add_accent_bar(s, L + Inches(0.3), top + Inches(0.3),
                       width=Inches(0.65), height=Inches(0.08), color=colors[i])
        add_text(s, L + Inches(0.3), top + Inches(0.5), card_w - Inches(0.6),
                 Inches(1.3), titles[i], size=20, bold=True, color=TEXT)
        bullet_list(s, L + Inches(0.3), top + Inches(2.4),
                    card_w - Inches(0.6), Inches(2.8), descs[i],
                    size=14, bullet_color=colors[i])
queue.append(slide_roadmap)


def slide_dataset(s):
    add_header(s, "Dataset · europe.csv")
    kpis = [
        ("Países", "28", ACCENT_T),
        ("Features", "7", ACCENT_P),
        ("Preprocesado", "Z-score", ACCENT_O),
        ("Var. explicada PC1", "46.1 %", ACCENT_T),
    ]
    kx = MARGIN_X; ky = CONTENT_TOP + Inches(0.1)
    kw = (CONTENT_W - Inches(0.6)) / 4
    for i, (lab, val, col) in enumerate(kpis):
        kpi_box(s, kx + (kw + Inches(0.2)) * i, ky, kw, Inches(1.65),
                lab, val, value_color=col, value_size=30)
    feats = [
        ["Feature", "Descripción"],
        ["Area", "Superficie (km²)"],
        ["GDP", "PBI per cápita (USD)"],
        ["Inflation", "Tasa de inflación anual (%)"],
        ["Life.expect", "Esperanza de vida (años)"],
        ["Military", "Gasto militar (% del PBI)"],
        ["Pop.growth", "Crecimiento poblacional (%)"],
        ["Unemployment", "Tasa de desempleo (%)"],
    ]
    add_table(s, MARGIN_X, ky + Inches(1.9), Inches(6.7), Inches(3.45),
              feats, font_size=13,
              col_widths=[Inches(1.9), Inches(4.8)])
    add_panel(s, MARGIN_X + Inches(6.95), ky + Inches(1.9),
              CONTENT_W - Inches(6.95), Inches(3.45), radius=True)
    add_text(s, MARGIN_X + Inches(7.20), ky + Inches(2.05),
             Inches(5), Inches(0.4),
             "Salidas del TP", size=14, bold=True, color=ACCENT_T)
    bullet_list(s, MARGIN_X + Inches(7.20), ky + Inches(2.55),
                Inches(5), Inches(3.0), [
        "Kohonen → country_map, u_matrix, hit_map",
        "Oja → loadings, country_scores, convergence",
        "Hopfield → recall (a), espúreo (b), ortogonalidad",
        "Hopfield → capacidad + escalado adaptativo",
        "Análisis cruzado → PCA vs Kohonen vs Oja",
    ], size=12)
queue.append(slide_dataset)


def slide_estandarizacion(s):
    add_header(s, "Estandarización · datos crudos vs Z-score")
    two_images(s, [
        (os.path.join(FIG, "pca_boxplot_raw.png"),
         "Datos crudos · Area y GDP dominan cualquier distancia"),
        (os.path.join(FIG, "pca_boxplot_std.png"),
         "Z-score · todas las features comparables"),
    ])
queue.append(slide_estandarizacion)


# ═══════════════════════ SECCIÓN 1.1 ═══════════════════════════════════════
def slide_section_11(s):
    set_bg(s, BG)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8),
                             SW, Inches(1.9))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.line.fill.background()
    add_accent_bar(s, Inches(0.55), Inches(3.05), Inches(0.10),
                   Inches(1.4), color=ACCENT_T)
    add_text(s, Inches(0.9), Inches(3.0), Inches(11), Inches(0.5),
             "EJERCICIO 1.1", size=16, bold=True, color=ACCENT_T)
    add_text(s, Inches(0.9), Inches(3.45), Inches(12), Inches(1.0),
             "Red de Kohonen (SOM)", size=44, bold=True, color=TEXT)
    add_text(s, Inches(0.9), Inches(4.25), Inches(12), Inches(0.5),
             "Agrupamiento topológico de países europeos",
             size=18, color=MUTED)
queue.append(slide_section_11)


def slide_kohonen_setup(s):
    add_header(s, "Setup · hiperparámetros SOM",
               section="EJERCICIO 1.1 · KOHONEN")
    data = [
        ["Parámetro", "Valor"],
        ["grid_rows × grid_cols", "5 × 5"],
        ["Inicialización pesos", "Gaussiana"],
        ["Tasa de aprendizaje (lr)", "0.5 → decaimiento exponencial"],
        ["Radio de vecindad", "3.0 → decaimiento exponencial"],
        ["Función de vecindad", "gaussiana"],
        ["Épocas", "1000"],
        ["Semilla", "42"],
        ["Distancia BMU", "Euclidiana sobre features estandarizadas"],
    ]
    table_w = Inches(8.5)
    add_table(s, MARGIN_X, CONTENT_TOP + Inches(0.2), table_w, Inches(5.0),
              data, col_widths=[Inches(3.2), table_w - Inches(3.2)],
              font_size=14)
    side_panel(s, MARGIN_X + table_w + IMG_PLUS_GAP,
               CONTENT_TOP + Inches(0.2),
               CONTENT_W - table_w - IMG_PLUS_GAP, Inches(5.0),
               heading="Salidas generadas",
               items=["country_map.png", "u_matrix.png", "hit_map.png"],
               item_size=13)
queue.append(slide_kohonen_setup)


def slide_country_map(s):
    add_header(s, "Country map · neurona BMU por país",
               section="EJERCICIO 1.1 · KOHONEN")
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Clusters observados",
                   items=["Nórdicos / ricos: NO",
                          "Europa del Este: SE",
                          "Mediterráneos: centro",
                          "Outlier: Ukraine",
                          "Lux. y Switzerland separados"])
    image_with_side(s, os.path.join(FIG, "kh_country_map.png"), side)
queue.append(slide_country_map)


def slide_u_matrix(s):
    add_header(s, "U-Matrix · distancias entre neuronas vecinas",
               section="EJERCICIO 1.1 · KOHONEN")
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Lectura",
                   items=["Crestas separan clusters",
                          "Valles agrupan países similares",
                          "Bordes nítidos contra Ukraine",
                          "Transición gradual en centro"])
    image_with_side(s, os.path.join(FIG, "kh_u_matrix.png"), side)
queue.append(slide_u_matrix)


def slide_hit_map(s):
    add_header(s, "Hit map · densidad por neurona",
               section="EJERCICIO 1.1 · KOHONEN")
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Distribución",
                   items=["Mayoría: 1–2 países / neurona",
                          "Pocas neuronas con 3+",
                          "Muchas vacías → buen spread",
                          "5×5 alcanza para 28"])
    image_with_side(s, os.path.join(FIG, "kh_hit_map.png"), side)
queue.append(slide_hit_map)


def slide_kohonen_vs_pca(s):
    add_header(s, "Comparación con PCA · PC1 sobre la grilla SOM",
               section="EJERCICIO 1.1 · KOHONEN ↔ PCA")
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Resultado",
                   items=["Gradiente continuo",
                          "Extremos PC1 en esquinas opuestas",
                          "SOM preserva la dirección de máxima varianza",
                          "Validación cruzada PCA ↔ Kohonen"])
    image_with_side(s, os.path.join(FIG, "cmp_pc1_gradient.png"), side)
queue.append(slide_kohonen_vs_pca)


def slide_kohonen_extremes(s):
    add_header(s, "Extremos de PC1 sobre el SOM",
               section="EJERCICIO 1.1 · KOHONEN ↔ PCA")
    def side(L, T, W, H):
        add_panel(s, L, T, W, H, radius=True)
        pad = Inches(0.22)
        add_text(s, L + pad, T + pad, W - 2 * pad, Inches(0.35),
                 "Top PC1", size=14, bold=True, color=GREEN)
        bullet_list(s, L + pad, T + pad + Inches(0.45),
                    W - 2 * pad, Inches(2.0),
                    ["Luxembourg", "Switzerland", "Norway"],
                    size=13, bullet_color=GREEN)
        add_text(s, L + pad, T + Inches(2.9), W - 2 * pad, Inches(0.35),
                 "Bottom PC1", size=14, bold=True, color=RED)
        bullet_list(s, L + pad, T + Inches(3.35),
                    W - 2 * pad, Inches(2.0),
                    ["Ukraine", "Bulgaria", "Estonia"],
                    size=13, bullet_color=RED)
    image_with_side(s, os.path.join(FIG, "cmp_kohonen_extremes.png"), side)
queue.append(slide_kohonen_extremes)


# ═══════════════════════ SECCIÓN 1.2 ═══════════════════════════════════════
def slide_section_12(s):
    set_bg(s, BG)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8),
                             SW, Inches(1.9))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.line.fill.background()
    add_accent_bar(s, Inches(0.55), Inches(3.05), Inches(0.10),
                   Inches(1.4), color=ACCENT_P)
    add_text(s, Inches(0.9), Inches(3.0), Inches(11), Inches(0.5),
             "EJERCICIO 1.2", size=16, bold=True, color=ACCENT_P)
    add_text(s, Inches(0.9), Inches(3.45), Inches(12), Inches(1.0),
             "Regla de Oja", size=44, bold=True, color=TEXT)
    add_text(s, Inches(0.9), Inches(4.25), Inches(12), Inches(0.5),
             "Primera componente principal por aprendizaje hebbiano",
             size=18, color=MUTED)
queue.append(slide_section_12)


def slide_oja_setup(s):
    add_header(s, "Setup · hiperparámetros Oja",
               section="EJERCICIO 1.2 · OJA", accent=ACCENT_P)
    data = [
        ["Parámetro", "Valor"],
        ["Input dim", "7 features"],
        ["Learning rate", "0.5"],
        ["Épocas", "1000"],
        ["Semilla", "42"],
        ["Inicialización", "Gaussiana, normalizada"],
        ["Comparación", "sklearn PCA(n_components=1)"],
    ]
    table_w = Inches(7.8)
    add_table(s, MARGIN_X, CONTENT_TOP + Inches(0.2), table_w, Inches(4.4),
              data, col_widths=[Inches(3.0), table_w - Inches(3.0)],
              font_size=14, header_bg=ACCENT_P)

    # KPI panel a la derecha
    panel_L = MARGIN_X + table_w + IMG_PLUS_GAP
    panel_W = CONTENT_W - table_w - IMG_PLUS_GAP
    add_panel(s, panel_L, CONTENT_TOP + Inches(0.2), panel_W, Inches(4.4),
              radius=True)
    add_text(s, panel_L + Inches(0.22), CONTENT_TOP + Inches(0.35),
             panel_W - Inches(0.4), Inches(0.4),
             "Validación numérica", size=14, bold=True, color=ACCENT_P)

    kpi_specs = [("cos(w_oja, w_sklearn)", "1.0000", ACCENT_T),
                 ("corr(scores)", "1.0000", ACCENT_T),
                 ("Var. explicada PC1", "46.10 %", ACCENT_P)]
    ky = CONTENT_TOP + Inches(0.9)
    for lab, val, col in kpi_specs:
        add_text(s, panel_L + Inches(0.22), ky, panel_W - Inches(0.4),
                 Inches(0.28), lab, size=10, color=MUTED)
        add_text(s, panel_L + Inches(0.22), ky + Inches(0.28),
                 panel_W - Inches(0.4), Inches(0.5),
                 val, size=22, bold=True, color=col)
        ky += Inches(1.10)

    add_text(s, MARGIN_X, CONTENT_TOP + Inches(4.85), CONTENT_W, Inches(0.4),
             "→ Oja converge al autovector dominante de la matriz de covarianza.",
             size=14, color=MUTED, align=PP_ALIGN.CENTER)
queue.append(slide_oja_setup)


def slide_oja_loadings(s):
    add_header(s, "Loadings de PC1 · Oja vs sklearn",
               section="EJERCICIO 1.2 · OJA", accent=ACCENT_P)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Interpretación",
                   items=["PC1 = eje de desarrollo",
                          "GDP, Life.expect, Pop.growth → positivo",
                          "Inflation, Unemployment → negativo",
                          "Barras Oja ≈ sklearn (overlap)"],
                   accent=ACCENT_P)
    image_with_side(s, os.path.join(FIG, "oja_loadings.png"), side)
queue.append(slide_oja_loadings)


def slide_oja_country_scores(s):
    add_header(s, "Países proyectados sobre PC1 · Oja",
               section="EJERCICIO 1.2 · OJA", accent=ACCENT_P)
    full_image(s, os.path.join(FIG, "oja_country_scores.png"))
queue.append(slide_oja_country_scores)


def slide_oja_convergence(s):
    add_header(s, "Convergencia hacia el autovector de sklearn",
               section="EJERCICIO 1.2 · OJA", accent=ACCENT_P)
    full_image(s, os.path.join(FIG, "oja_convergence.png"),
               caption="Distancia y similitud coseno época a época")
queue.append(slide_oja_convergence)


def slide_oja_vs_sklearn_scatter(s):
    add_header(s, "Scores Oja vs sklearn · 28 países",
               section="EJERCICIO 1.2 · OJA", accent=ACCENT_P)
    def side(L, T, W, H):
        add_panel(s, L, T, W, H, radius=True)
        pad = Inches(0.22)
        add_text(s, L + pad, T + pad, W - 2 * pad, Inches(0.4),
                 "Equivalencia exacta", size=14, bold=True, color=ACCENT_T)
        y = T + pad + Inches(0.6)
        for lab, val in [("corr(scores)", "1.0000"),
                         ("cos(w_oja, w_sklearn)", "1.0000")]:
            add_text(s, L + pad, y, W - 2 * pad, Inches(0.28),
                     lab, size=11, color=MUTED)
            add_text(s, L + pad, y + Inches(0.28), W - 2 * pad, Inches(0.5),
                     val, size=22, bold=True, color=ACCENT_T)
            y += Inches(1.0)
        add_text(s, L + pad, y + Inches(0.2), W - 2 * pad, Inches(1.5),
                 "Cada punto cae sobre y = x: Oja reproduce la proyección de sklearn sin diferencia detectable.",
                 size=11, color=MUTED)
    image_with_side(s, os.path.join(FIG, "oja_vs_sklearn_scatter.png"), side)
queue.append(slide_oja_vs_sklearn_scatter)


# ═══════════════════════ SECCIÓN 2 ═════════════════════════════════════════
def slide_section_2(s):
    set_bg(s, BG)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.8),
                             SW, Inches(1.9))
    box.fill.solid(); box.fill.fore_color.rgb = PANEL; box.line.fill.background()
    add_accent_bar(s, Inches(0.55), Inches(3.05), Inches(0.10),
                   Inches(1.4), color=ACCENT_O)
    add_text(s, Inches(0.9), Inches(3.0), Inches(11), Inches(0.5),
             "EJERCICIO 2", size=16, bold=True, color=ACCENT_O)
    add_text(s, Inches(0.9), Inches(3.45), Inches(12), Inches(1.0),
             "Hopfield · memoria asociativa", size=42, bold=True, color=TEXT)
    add_text(s, Inches(0.9), Inches(4.25), Inches(12), Inches(0.5),
             "Patrones de letras 5×5 con ruido y estados espúreos",
             size=18, color=MUTED)
queue.append(slide_section_2)


def slide_hopfield_setup(s):
    add_header(s, "Setup · subset elegido y configuración",
               section="EJERCICIO 2 · HOPFIELD", accent=ACCENT_O)
    data = [
        ["Parámetro", "Valor"],
        ["Patrones almacenados", "G, R, T, V"],
        ["Tamaño de patrón", "5 × 5 (N = 25)"],
        ["Regla", "Hebb, W_ii = 0"],
        ["Modo de actualización", "Síncrono"],
        ["Ruido (parte a)", "15 % de bits invertidos"],
        ["Ruido alto (parte b)", "40 %"],
        ["Max steps", "50"],
        ["Semilla", "42"],
    ]
    table_w = Inches(8.0)
    add_table(s, MARGIN_X, CONTENT_TOP + Inches(0.2), table_w, Inches(5.2),
              data, col_widths=[Inches(2.8), table_w - Inches(2.8)],
              font_size=14, header_bg=ACCENT_O)
    side_panel(s, MARGIN_X + table_w + IMG_PLUS_GAP,
               CONTENT_TOP + Inches(0.2),
               CONTENT_W - table_w - IMG_PLUS_GAP, Inches(5.2),
               heading="¿Por qué GRTV?",
               items=["Barrido C(26,4) = 14 950 combinaciones",
                      "Métrica: max |⟨ξᵢ, ξⱼ⟩|",
                      "GRTV → max = 1 (casi ortogonal)",
                      "Óptimo para N = 25"],
               accent=ACCENT_O, item_size=12)
queue.append(slide_hopfield_setup)


# ── (a) recall — un slide por letra ─────────────────────────────────────────
def make_recall_slide(letter, file_key):
    def _builder(s):
        add_header(s, f"(a) Recall paso a paso · letra {letter}",
                   section="EJERCICIO 2 · HOPFIELD · PARTE (a)",
                   accent=ACCENT_O)
        full_image(s, os.path.join(FIG, f"hop_recall_a_{file_key}.png"),
                   caption=f"Iteraciones síncronas hasta converger a {letter}")
    return _builder


for L in ["G", "R", "T", "V"]:
    queue.append(make_recall_slide(L, L))


# ── (a) energía — 2x2 compacto ──────────────────────────────────────────────
def slide_energy_a(s):
    add_header(s, "(a) Energía vs iteración · 4 letras",
               section="EJERCICIO 2 · HOPFIELD · PARTE (a)",
               accent=ACCENT_O)
    grid_2x2(s, [
        (os.path.join(FIG, "hop_energy_a_G.png"), "G"),
        (os.path.join(FIG, "hop_energy_a_R.png"), "R"),
        (os.path.join(FIG, "hop_energy_a_T.png"), "T"),
        (os.path.join(FIG, "hop_energy_a_V.png"), "V"),
    ])
queue.append(slide_energy_a)


def slide_recall_metrics(s):
    add_header(s, "(a) Métricas de recall del subset GRTV",
               section="EJERCICIO 2 · HOPFIELD · PARTE (a)",
               accent=ACCENT_O)
    kpis = [
        ("Letras almacenadas", "4", ACCENT_T),
        ("Recall accuracy (15% ruido)", "100 %", GREEN),
        ("Hamming final promedio", "0.0", ACCENT_T),
        ("Estados espúreos en (a)", "0", GREEN),
    ]
    kx = MARGIN_X; ky = CONTENT_TOP + Inches(0.3)
    kw = (CONTENT_W - Inches(0.6)) / 4
    for i, (lab, val, col) in enumerate(kpis):
        kpi_box(s, kx + (kw + Inches(0.2)) * i, ky, kw, Inches(2.1),
                lab, val, value_color=col, value_size=36)
    add_panel(s, MARGIN_X, ky + Inches(2.4), CONTENT_W, Inches(2.85),
              radius=True)
    add_text(s, MARGIN_X + Inches(0.3), ky + Inches(2.6),
             CONTENT_W - Inches(0.6), Inches(0.4),
             "Lectura", size=15, bold=True, color=ACCENT_O)
    bullet_list(s, MARGIN_X + Inches(0.3), ky + Inches(3.1),
                CONTENT_W - Inches(0.6), Inches(2.2), [
        "Con un subset ortogonal y ruido bajo, Hopfield recupera el 100 % en ≤ 5 iteraciones",
        "Energía monótonamente no creciente bajo update síncrono (sin ciclos detectados)",
        "Punto de partida limpio para estudiar capacidad y crosstalk",
    ], size=14)
queue.append(slide_recall_metrics)


# ── (b) espúreo ─────────────────────────────────────────────────────────────
def slide_spurious_recall(s):
    add_header(s, "(b) Patrón muy ruidoso · trayectoria de recall",
               section="EJERCICIO 2 · HOPFIELD · PARTE (b)",
               accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_recall_b.png"),
               caption="40 % de ruido — converge a un estado fijo que no es ninguna de las 4 letras almacenadas")
queue.append(slide_spurious_recall)


def slide_spurious_energy(s):
    add_header(s, "(b) Energía del estado espúreo",
               section="EJERCICIO 2 · HOPFIELD · PARTE (b)",
               accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_energy_b.png"),
               caption="La red cae en un mínimo de energía que es punto fijo pero no es un patrón almacenado")
queue.append(slide_spurious_energy)


# ── Extras Hopfield ────────────────────────────────────────────────────────
def slide_abecedario(s):
    add_header(s, "Abecedario completo · 26 patrones 5×5",
               section="EJERCICIO 2 · EXTRAS", accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_abecedario.png"))
queue.append(slide_abecedario)


def slide_ortho_heatmap(s):
    add_header(s, "Ortogonalidad · matriz 26 × 26 de |⟨ξᵢ, ξⱼ⟩|",
               section="EJERCICIO 2 · EXTRAS · ORTOGONALIDAD",
               accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Lectura",
                   items=["Diagonal = 25 (auto-correlación)",
                          "Pares ≥ 15 (E↔F, C↔G)",
                          "Hopfield falla con patrones correlacionados",
                          "Motiva el subset óptimo"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_dot_heatmap.png"), side)
queue.append(slide_ortho_heatmap)


def slide_distribution_k4(s):
    add_header(s, "Distribución de max|⟨·,·⟩| sobre 14 950 subsets",
               section="EJERCICIO 2 · EXTRAS · ORTOGONALIDAD",
               accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Resultados",
                   items=["Mínimo absoluto = 1 (GRTV)",
                          "Cola izquierda corta",
                          "Mayoría con max-dot ≥ 5",
                          "Elección importa más que el azar"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_distribution_k4.png"), side)
queue.append(slide_distribution_k4)


def slide_top_bottom(s):
    add_header(s, "Top vs bottom · mejores y peores subsets de 4 letras",
               section="EJERCICIO 2 · EXTRAS · ORTOGONALIDAD",
               accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_top_bottom_k4.png"),
               caption="Arriba: los más ortogonales · Abajo: los menos ortogonales")
queue.append(slide_top_bottom)


def slide_crosstalk(s):
    add_header(s, "Crosstalk del subset GRTV · ξᵢ · ξⱼ / N",
               section="EJERCICIO 2 · EXTRAS", accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Métricas",
                   items=["Diagonal = 1",
                          "Off-diagonal ≤ 0.04",
                          "Crosstalk ≈ 0",
                          "Justifica 100 % de recall en (a)"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_crosstalk.png"), side)
queue.append(slide_crosstalk)


def slide_recovery(s):
    add_header(s, "Recovery rate vs ruido · una curva por letra",
               section="EJERCICIO 2 · EXTRAS", accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_recovery.png"),
               caption="Caída abrupta cerca de 40-45 %: la cuenca de atracción tiene radio finito")
queue.append(slide_recovery)


# ── Capacidad — un slide por métrica ───────────────────────────────────────
def slide_capacity_accuracy(s):
    add_header(s, "Capacidad · recall accuracy vs N de patrones",
               section="EJERCICIO 2 · EXTRAS · CAPACIDAD", accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Lectura",
                   items=["Caída marcada al pasar 3 patrones",
                          "Coincide con el límite teórico",
                          "p ≈ 0.138 · N = 3.45 para N = 25",
                          "Misma tendencia con/sin ruido"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_accuracy_vs_n.png"), side)
queue.append(slide_capacity_accuracy)


def slide_capacity_spurious(s):
    add_header(s, "Capacidad · tasa de estados espúreos vs N",
               section="EJERCICIO 2 · EXTRAS · CAPACIDAD", accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Lectura",
                   items=["Crece monótonamente con N",
                          "≥ 50 % de las consultas caen en espúreos a partir de N=5",
                          "La red prefiere mínimos compuestos"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_spurious_vs_n.png"), side)
queue.append(slide_capacity_spurious)


def slide_capacity_hamming(s):
    add_header(s, "Capacidad · Hamming promedio al patrón original",
               section="EJERCICIO 2 · EXTRAS · CAPACIDAD", accent=ACCENT_O)
    def side(L, T, W, H):
        side_panel(s, L, T, W, H, heading="Lectura",
                   items=["Hamming ≈ 0 cuando recall ok",
                          "Distancia explota al saturar",
                          "Plateau cerca de N/2 bits errados"],
                   accent=ACCENT_O)
    image_with_side(s, os.path.join(FIG, "hop_hamming_vs_n.png"), side)
queue.append(slide_capacity_hamming)


# ── Escalado ────────────────────────────────────────────────────────────────
def slide_scaling_table(s):
    add_header(s, "Escalado adaptativo · np.kron(5×5 → 5k×5k)",
               section="EJERCICIO 2 · EXTRAS · ESCALADO", accent=ACCENT_O)
    data = [
        ["Patrones (p)", "k", "N = (5k)²", "p / N", "≤ 0.138 N?"],
        ["≤ 3", "1", "25", "≤ 0.12", "✓"],
        ["4 – 13", "2", "100", "≤ 0.13", "✓"],
        ["14 – 31", "3", "225", "≤ 0.116", "✓"],
        ["26 (alfabeto)", "3", "225", "0.116", "✓"],
    ]
    table_w = Inches(8.6)
    add_table(s, MARGIN_X, CONTENT_TOP + Inches(0.3), table_w, Inches(3.6),
              data, font_size=15, header_bg=ACCENT_O,
              col_widths=[Inches(2.0), Inches(0.8), Inches(1.7),
                          Inches(1.8), Inches(2.3)])
    side_panel(s, MARGIN_X + table_w + IMG_PLUS_GAP,
               CONTENT_TOP + Inches(0.3),
               CONTENT_W - table_w - IMG_PLUS_GAP, Inches(3.6),
               heading="Idea clave",
               items=["Cada bit → bloque k×k vía Kronecker",
                      "Crece N sin alterar el patrón",
                      "Respeta p ≤ 0.138 · N",
                      "Con k=3 entran las 26 letras"],
               accent=ACCENT_O, item_size=12)
    add_panel(s, MARGIN_X, CONTENT_TOP + Inches(4.1),
              CONTENT_W, Inches(1.6), radius=True)
    add_text(s, MARGIN_X + Inches(0.3), CONTENT_TOP + Inches(4.25),
             CONTENT_W - Inches(0.6), Inches(0.4),
             "Por qué importa", size=14, bold=True, color=ACCENT_T)
    add_text(s, MARGIN_X + Inches(0.3), CONTENT_TOP + Inches(4.65),
             CONTENT_W - Inches(0.6), Inches(0.9),
             "Con N = 25 y p = 26 la red colapsa (p/N ≈ 1.04). El escalado "
             "adaptativo desbloquea el almacenamiento del alfabeto entero sin "
             "tocar la regla de Hebb.",
             size=13, color=MUTED)
queue.append(slide_scaling_table)


def slide_fixed_vs_adaptive(s):
    add_header(s, "Capacidad fija (N=25) vs adaptativa",
               section="EJERCICIO 2 · EXTRAS · ESCALADO", accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_fixed_vs_adaptive.png"),
               caption="Línea teórica 0.138 · N marcada como referencia · adaptativo mantiene recall alto")
queue.append(slide_fixed_vs_adaptive)


def slide_alphabet_crosstalk(s):
    add_header(s, "Alfabeto completo (k=3) · crosstalk",
               section="EJERCICIO 2 · EXTRAS · ESCALADO", accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_crosstalk_alphabet.png"),
               caption="Crosstalk de las 26 letras escaladas — pares parecidos (A/H, C/G/O, E/F) se confunden por correlación, no por capacidad")
queue.append(slide_alphabet_crosstalk)


def slide_alphabet_recovery(s):
    add_header(s, "Alfabeto completo (k=3) · recovery vs ruido por letra",
               section="EJERCICIO 2 · EXTRAS · ESCALADO", accent=ACCENT_O)
    full_image(s, os.path.join(FIG, "hop_recovery_alphabet.png"),
               caption="Las letras menos correlacionadas mantienen recall alto hasta ~40 % de ruido")
queue.append(slide_alphabet_recovery)


# ═══════════════════════ CIERRE ═════════════════════════════════════════════
def slide_conclusiones(s):
    add_header(s, "Conclusiones")
    card_w = Inches(4.05); card_h = Inches(5.4)
    gap = Inches(0.20); start = MARGIN_X; top = CONTENT_TOP + Inches(0.10)
    titles = ["Kohonen", "Oja", "Hopfield"]
    points = [
        ["Mapa topológico estable",
         "Clusters geo-económicos claros",
         "Extremos PC1 en esquinas opuestas",
         "Validación cruzada con PCA"],
        ["Converge al autovector dominante",
         "cos = 1.0000 con sklearn",
         "corr(scores) = 1.0000",
         "PC1 ≈ eje de desarrollo socio-económico"],
        ["100 % recall con subset ortogonal (GRTV)",
         "Crosstalk domina antes que capacidad",
         "Quiebre a p ≈ 0.138 · N (= 3.45)",
         "Escalado np.kron habilita las 26 letras"],
    ]
    colors = [ACCENT_T, ACCENT_P, ACCENT_O]
    for i in range(3):
        L = start + (card_w + gap) * i
        add_panel(s, L, top, card_w, card_h, radius=True)
        add_accent_bar(s, L + Inches(0.3), top + Inches(0.3),
                       width=Inches(0.65), height=Inches(0.08), color=colors[i])
        add_text(s, L + Inches(0.3), top + Inches(0.5), card_w - Inches(0.6),
                 Inches(0.7), titles[i], size=22, bold=True, color=TEXT)
        bullet_list(s, L + Inches(0.3), top + Inches(1.6),
                    card_w - Inches(0.6), Inches(3.6), points[i],
                    size=13, bullet_color=colors[i])
queue.append(slide_conclusiones)


def slide_entrega(s):
    add_header(s, "Entrega · repo + commit")
    add_panel(s, MARGIN_X, CONTENT_TOP + Inches(0.2),
              CONTENT_W, Inches(4.4), radius=True)
    add_text(s, MARGIN_X + Inches(0.3), CONTENT_TOP + Inches(0.5),
             CONTENT_W - Inches(0.6), Inches(0.5),
             "Reproducibilidad", size=18, bold=True, color=ACCENT_T)
    bullet_list(s, MARGIN_X + Inches(0.3), CONTENT_TOP + Inches(1.1),
                CONTENT_W - Inches(0.6), Inches(3.3), [
        "Repositorio: SIA-G02 · branch tp4 · commit 28217db",
        "Configs: configs/{kohonen_europe,oja_europe,hopfield}.json",
        "Reproducción end-to-end:  make all",
        "Tests:  uv run pytest    · cubre Hopfield/SOM/Oja/PCA",
        "Salidas: output/ + kohonen/output/ + pca_test/plots/",
    ], size=14)
    add_text(s, MARGIN_X, CONTENT_TOP + Inches(5.0), CONTENT_W, Inches(0.5),
             "Gracias.", size=22, bold=True, color=ACCENT_T,
             align=PP_ALIGN.CENTER)
queue.append(slide_entrega)


# ─── Render ─────────────────────────────────────────────────────────────────
TOTAL = len(queue)
for i, build in enumerate(queue, 1):
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    build(s)
    add_footer(s, i, TOTAL)


prs.save(OUT)
print(f"OK — {TOTAL} slides → {OUT}")
